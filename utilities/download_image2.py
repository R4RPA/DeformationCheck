import win32com.client
import PIL
from PIL import ImageGrab, Image
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.util import Pt
from utilities import variables
import ctypes
from time import sleep
import shutil
import json

class DownloadImage:
    def __init__(self, data):
        """Initialize the class with the paths of the image, input files, and output file."""
        self.image_path = data["image_path"]
        self.ip_file_1 = data["ip_file_1"]
        self.ip_file_2 = data["ip_file_2"]
        self.op_file = data["op_file"]
        self.sv_file = data["sv_file"]

    def extract_images_from_excel(self):
        """This function extracts a graph from the input excel file and saves it into the specified PNG image path (overwrites the given PNG image)"""
        print("Process input file 1 : " + self.ip_file_1)
        self.save_excel_chart_as_jpg(self.ip_file_1)
        print("Process input file 2 : " + self.ip_file_2)
        self.save_excel_chart_as_jpg(self.ip_file_2)
    
    def crete_ppt(self):
        presentation = Presentation(self.op_file)
        presentation.save(self.sv_file)
        
    def copy_downloaded_images_to_ppt(self):
        """This function inserts downloaded imaged into respective ppt slide with reference to definition available in varaiables_dict"""
        # Open the PowerPoint presentation
        self.presentation = Presentation(self.sv_file)

        for img_file_name, val in variables.image_filename_slide_dict.items():
            slide_idx = val["slide_index"]
            print("Inserting " + img_file_name + " into slide " + str(slide_idx) + " of file " + str(self.sv_file))
            image_file = str(self.image_path + img_file_name)
            dmns = variables.dimensions_dict[val["image_type"]]
            self.copy_image(image_file, slide_idx, dmns["image_height"], dmns["image_width"],
                            dmns["horizontal_pos"], dmns["vertical_pos"])
        self.presentation.save(self.sv_file)

    def save_excel_chart_as_jpg(self, filename):
        """This function to extact shapes and charts from excel and save in Images folder"""
        print("Processing the excel charts and saving it to the images folder in the current directory")
        self.killexcel()
        # Open the excel application using win32com
        o = win32com.client.Dispatch("Excel.Application")
        # Disable alerts and visibility to the user
        o.Visible = 0
        o.DisplayAlerts = 0
        # Open workbook
        wb = o.Workbooks.Open(filename, None, False)
        sleep(5)
        table_reference_dict = {"Summary": variables.summary_reference_dict,
                                "Graph":  variables.judgment_reference_dict}

        for m, sheet in enumerate(o.Sheets):
            sheet.Cells.ClearComments()
            if sheet.Name in ["Summary", "Graph"]:
                reference_dict = table_reference_dict[sheet.Name]['params']
                output_filename = table_reference_dict[sheet.Name]['file_name']
                self.create_data_files(sheet, reference_dict, sheet.Name, output_filename)
                if sheet.Name == "Summary":
                    reference_dict = variables.radius_reference_dict['params']
                    output_filename = variables.radius_reference_dict['file_name']
                    self.create_data_files(sheet, reference_dict, sheet.Name, output_filename)
                    

            for n, shape in enumerate(sheet.Shapes):
                new_img_path = f"{self.image_path}{sheet.Name}_{n}"
                new_img_path = new_img_path.replace(" ", "").replace(".", "") + ".png"
                if os.path.exists(new_img_path):
                    new_img_path = f"{self.image_path}{sheet.Name}_{n}{n}"
                    new_img_path = new_img_path.replace(" ", "").replace(".", "") + ".png"
                print('shape', shape.Name, 'new_img_path', new_img_path)
                try_copy_image = True
                wait_time = 0.5
                while try_copy_image:
                    try:
                        self.clear_clip_board()
                        shape.Copy()
                        sleep(wait_time)
                        image = ImageGrab.grabclipboard().convert('RGB')
                        image.save(new_img_path)
                        try_copy_image = False
                    except:
                        wait_time += 1
                        print('error - retrying, wait_time =', wait_time)
        o.Quit()
            
    def create_data_files(self, ws, reference_dict, sheet_name, output_filename):
        """this function to create summary and jusgment table json files"""
        data = {}
        for component, stages in reference_dict.items():
            print('component', component, 'stages', stages)
            for stage, cell_range in stages.items():
                try:
                    range_obj = ws.Range(cell_range)
                    range_data = []
                    for row in range_obj.Rows:
                        row_data = []
                        for cell in row.Cells:
                            if cell.Value:
                                row_data.append(cell.Text)
                            else:
                                row_data.append('')
                        if sheet_name == 'Summary':
                            range_data.append(row_data)
                        else:
                            range_data.append([row_data[0]])
                    if component not in data:
                        data[component] = {}
                    data[component][stage] = range_data
                except Exception as e:
                    print(f"Error occurred while trying to access cell range {cell_range}: {e}")
        output_file = os.path.join(self.image_path, output_filename)
        with open(output_file, 'w') as f:
            json.dump(data, f)
        
    def copy_image(self, image_path, slide_index, img_h, img_w, horizontal_pos, vertical_pos):
        """Inserts an image into the specified slide in the PowerPoint presentation."""
        slide = self.presentation.slides[slide_index]

        # Setting Image size
        image_width = Cm(img_w)
        image_height = Cm(img_h)

        # Setting position of image in the slide
        left = Cm(horizontal_pos)
        top = Cm(vertical_pos)

        slide.shapes.add_picture(image_path, left, top, image_width, image_height)
        
    
    def update_pptx_tables(self):
        """this function to update summary and jusgment table in the slides"""
        summary_file = variables.summary_reference_dict['file_name']
        judgement_file = variables.judgment_reference_dict['file_name']
        radius_file = variables.radius_reference_dict['file_name']
        summary_file = os.path.join(self.image_path, summary_file)
        judgement_file = os.path.join(self.image_path, judgement_file)
        radius_file = os.path.join(self.image_path, radius_file)
        with open(summary_file, 'r') as f:
            summary_data = json.load(f)
        with open(judgement_file, 'r') as f:
            judgement_data = json.load(f)
        with open(radius_file, 'r') as f:
            radius_data = json.load(f)
    
        prs = Presentation(self.sv_file)
        slide_indices = variables.slide_indices
    
        for component in summary_data:
            for i, stage in enumerate(summary_data[component]):
                slide = prs.slides[slide_indices[component][i]]
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        # If the table has only 1 column, treat it as a judgment table
                        if len(table.columns) == 1:
                            judgement_values = judgement_data[component][stage]
                            for row in range(len(judgement_values)):
                                cell = table.cell(row, 0)
                                cell.text = judgement_values[row][0]
                        # If the table has more than 1 column, treat it as a summary table
                        elif len(table.columns) > 1:
                            summary_values = summary_data[component][stage]
                            for row in range(len(summary_values)):
                                for col in range(len(summary_values[row])):
                                    cell = table.cell(row, col)
                                    cell.text = str(summary_values[row][col])
                        font_size_pt = Pt(12)
                        for cell in table.iter_cells():
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = font_size_pt
                                    
        for component in radius_data:
            for i, stage in enumerate(radius_data[component]):
                slide = prs.slides[slide_indices[component][i]]
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        if len(table.columns) > 1:
                            radius_values = radius_data[component][stage]
                            for row in range(len(radius_values)):
                                for col in range(len(radius_values[row])):
                                    try:
                                        cell = table.cell(row, col)
                                        cell.text = str(radius_values[row][col])
                                    except:
                                        pass
                        font_size_pt = Pt(12)
                        for cell in table.iter_cells():
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = font_size_pt
                                    
        prs.save(self.sv_file)
        
    def create_folder_in_current_directory(self):
        """This function will create Images folder if does not exist"""
        folder_path = os.path.join(self.image_path)
        try:
            os.mkdir(folder_path)
            print(f"Folder '{self.image_path}' created successfully.")
        except FileExistsError:
            print(f"Folder '{self.image_path}' already exists.")

    def delete_files_in_folder(self):
        """This function will clear Images folder contents"""
        print('Delete image file -Start')
        shutil.rmtree(self.image_path)
        print('Deleted image files')

    def clear_clip_board(self):
        """This function will clear clipboard"""
        loop_count = 0
        while True:
            loop_count += 1
            try:
                ctypes.windll.user32.OpenClipboard(None)
                ctypes.windll.user32.EmptyClipboard()
                ctypes.windll.user32.CloseClipboard()
                break
            except:
                sleep(5)
                if loop_count > 5:
                    break

    def killexcel(self):
        """This function will Kill/Force Close Excel application"""
        try:
            os.system('taskkill /IM EXCEL.exe /T /F')
        except:
            pass
        sleep(5)

